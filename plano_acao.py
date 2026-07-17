import streamlit as st
from supabase import create_client
from datetime import datetime, timezone, timedelta, date
import pandas as pd
from io import BytesIO


SUPABASE_URL = "https://yjvhfhodhlxgprpxelxs.supabase.co"
SUPABASE_KEY = st.secrets["SUPABASE_KEY"]

BRASILIA = timezone(timedelta(hours=-3))

def agora_brasilia():
    return datetime.now(BRASILIA)

RESPONSAVEIS = sorted([
    "Antonio", "Aline", "Andreia", "Amanda Medeiros", "Borges", "Cícera", "Cristiane",
    "Dani Vieira", "Diego", "Diego/Antonio", "Ednalva", "Eunice",
    "Fauzi", "Flávia", "Galvão", "Giselle", "Guilherme Carrasco", "Hebert",
    "Ieda", "Irlanilson", "Izabella", "Janaina Cruz", "Janaina/OGS", "Kleiton",
    "Lia Mara", "Luis Fernando", "Marketing", "Natalia", "Núbia",
    "Rogerio Galvão", "Roldão", "Stephanie", "Tânia", "TI", "Viviane"
])

STATUS_OPCOES = ["Em andamento", "Concluído", "Atrasado", "Cancelado"]
STATUS_ICONS = {"Concluído": "✅", "Atrasado": "⚠️", "Em andamento": "🔄", "Cancelado": "❌"}
TIPOS = ["Feedback/Treinamento", "Sistema", "Melhoria de Processo", "Projetos", "Outros", "Antigo"]
CORES_TIPO = {
    "Feedback/Treinamento": "#0055AA",
    "Sistema": "#CC0000",
    "Melhoria de Processo": "#1A7A1A",
    "Projetos": "#CC7700",
    "Outros": "#666666",
    "Antigo": "#999999",
}

st.set_page_config(page_title="Plano de Ação | Portabilidade", layout="wide")

st.markdown("""
<style>
.header-box { background-color: #cc0000; padding: 24px 32px; border-radius: 8px; margin-bottom: 24px; }
.header-box h1 { color: white; margin: 0; font-size: 2rem; }
.header-box p { color: rgba(255,255,255,0.85); margin: 6px 0 0 0; font-size: 0.9rem; }
</style>
""", unsafe_allow_html=True)

@st.cache_resource
def get_supabase():
    return create_client(SUPABASE_URL, SUPABASE_KEY)

def carregar_dados():
    sb = get_supabase()
    result = sb.table("plano_acao").select("*").order("numero").execute()
    return pd.DataFrame(result.data)

def carregar_historico(acao_id):
    sb = get_supabase()
    result = sb.table("plano_acao_historico").select("*").eq("acao_id", acao_id).order("atualizado_em", desc=True).execute()
    return pd.DataFrame(result.data)

def salvar_historico(acao_id, numero, status_anterior, status_novo, comentario, novo_prazo, responsavel):
    sb = get_supabase()
    sb.table("plano_acao_historico").insert({
        "acao_id": int(acao_id), "numero": int(numero),
        "status_anterior": status_anterior, "status_novo": status_novo,
        "comentario": comentario, "novo_prazo": novo_prazo,
        "atualizado_por": responsavel, "atualizado_em": agora_brasilia().isoformat(),
    }).execute()

def atualizar_registro(id_, status, comentario, responsavel, novo_prazo=None, tipo=None):
    sb = get_supabase()
    update = {
        "status": status, "comentario": comentario or None,
        "atualizado_por": responsavel, "atualizado_em": agora_brasilia().isoformat(),
    }
    if novo_prazo: update["prazo"] = novo_prazo
    if tipo: update["tipo"] = tipo
    sb.table("plano_acao").update(update).eq("id", id_).execute()

def cadastrar_acao(dados):
    get_supabase().table("plano_acao").insert(dados).execute()

def formatar_data(dt_str):
    if not dt_str: return ""
    try:
        dt = datetime.fromisoformat(str(dt_str))
        if dt.tzinfo is None: dt = dt.replace(tzinfo=timezone.utc)
        return dt.astimezone(BRASILIA).strftime("%d/%m/%Y %H:%M")
    except: return str(dt_str)

def limpar(valor):
    if valor is None: return None
    s = str(valor).strip()
    return None if s.lower() in ('nan', 'none', '') else s

def dias_para_prazo(prazo_str):
    if not prazo_str: return None
    try:
        prazo = datetime.strptime(str(prazo_str)[:10], "%Y-%m-%d").date()
        return (prazo - agora_brasilia().date()).days
    except: return None

def gerar_ppt_ata(df_ata, titulo_ata):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    RED = RGBColor(0xCC, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x1A, 0x1A, 0x1A)
    GRAY = RGBColor(0x66, 0x66, 0x66)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_rect(slide, x, y, w, h, color):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_text(slide, text, x, y, w, h, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
        txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    # SLIDE 1 — CAPA
    s1 = prs.slides.add_slide(blank)
    add_rect(s1, 0, 0, 13.33, 3.5, RED)
    add_text(s1, "📋 Plano de Ação — Portabilidade", 0.5, 0.6, 12, 1, size=36, bold=True, color=WHITE)
    add_text(s1, titulo_ata, 0.5, 1.7, 12, 0.6, size=18, color=WHITE)
    add_text(s1, f"Data: {agora_brasilia().strftime('%d/%m/%Y %H:%M')} | Total de ações: {len(df_ata)}", 0.5, 2.4, 12, 0.5, size=14, color=RGBColor(0xFF,0xCC,0xCC))

    # Métricas na capa
    status_counts = df_ata['status'].value_counts()
    metricas = [
        ("🔄 Em andamento", status_counts.get("Em andamento", 0), RGBColor(0x00,0x55,0xAA)),
        ("⚠️ Atrasadas", status_counts.get("Atrasado", 0), RGBColor(0xCC,0x77,0x00)),
        ("✅ Concluídas", status_counts.get("Concluído", 0), RGBColor(0x1A,0x7A,0x1A)),
    ]
    for i, (label, val, cor) in enumerate(metricas):
        x = 1.5 + i * 3.5
        add_rect(s1, x, 4.0, 3.0, 1.5, RGBColor(0xF5,0xF5,0xF5))
        add_text(s1, str(val), x, 4.1, 3.0, 0.8, size=36, bold=True, color=cor, align=PP_ALIGN.CENTER)
        add_text(s1, label, x, 4.9, 3.0, 0.4, size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # SLIDES POR STATUS
    for status in ["Em andamento", "Atrasado", "Concluído"]:
        df_s = df_ata[df_ata["status"] == status]
        if df_s.empty: continue

        icon = STATUS_ICONS.get(status, "")

        # Dividir em grupos de 5 por slide
        chunks = [df_s.iloc[i:i+5] for i in range(0, len(df_s), 5)]
        for chunk_idx, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(blank)
            add_rect(slide, 0, 0, 13.33, 1.1, RED)
            add_text(slide, f"{icon} {status}", 0.4, 0.15, 12, 0.75, size=26, bold=True, color=WHITE)
            if len(chunks) > 1:
                add_text(slide, f"({chunk_idx+1}/{len(chunks)})", 11.0, 0.25, 2, 0.5, size=12, color=WHITE)

            y = 1.25
            for _, row in chunk.iterrows():
                prazo = limpar(row.get('prazo')) or 'Sem prazo'
                tipo = limpar(row.get('tipo')) or '-'
                responsavel = row.get('responsavel', '')
                problema = str(row.get('problema_identificado', ''))[:120]
                plano = str(row.get('plano_de_acao', ''))[:100]
                comentario = limpar(row.get('comentario'))

                # Card da ação
                add_rect(slide, 0.3, y, 12.7, 1.0, RGBColor(0xF9,0xF9,0xF9))
                add_text(slide, f"#{int(row['numero'])} — {responsavel} | Prazo: {prazo} | Tipo: {tipo}",
                         0.45, y + 0.05, 12.3, 0.3, size=11, bold=True, color=RED)
                add_text(slide, f"📌 {problema}", 0.45, y + 0.35, 12.3, 0.28, size=10, color=DARK)
                if comentario:
                    add_text(slide, f"💬 {comentario}", 0.45, y + 0.65, 12.3, 0.25, size=9, color=GRAY)
                else:
                    add_text(slide, f"🎯 {plano}", 0.45, y + 0.65, 12.3, 0.25, size=9, color=GRAY)
                y += 1.1

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ============================
# SIDEBAR
# ============================
with st.sidebar:
    st.header("👤 Identificação")
    usuario = st.selectbox("Seu nome", ["Selecione..."] + RESPONSAVEIS)

# ============================
# DADOS
# ============================
df_all = carregar_dados()
total = len(df_all)
concluidas = len(df_all[df_all["status"] == "Concluído"])
atrasadas = len(df_all[df_all["status"] == "Atrasado"])
em_andamento = len(df_all[df_all["status"] == "Em andamento"])
taxa = f"{round(concluidas / total * 100, 1)}%" if total > 0 else "0%"
agora_str = agora_brasilia().strftime("%d/%m/%Y %H:%M")

# HEADER
st.markdown(f"""
<div class="header-box">
    <h1>📋 Plano de Ação — Portabilidade</h1>
    <p>Atualizado automaticamente · {agora_str}</p>
</div>
""", unsafe_allow_html=True)

# MÉTRICAS
col1, col2, col3, col4, col5 = st.columns(5)
col1.metric("Total de Ações", total)
col2.metric("✅ Concluídas", concluidas)
col3.metric("⚠️ Atrasadas", atrasadas)
col4.metric("🔄 Em andamento", em_andamento)
col5.metric("📈 Taxa de Conclusão", taxa)

# Alerta vencimento
df_vence = df_all[(df_all["status"].isin(["Em andamento", "Atrasado"])) & (df_all["prazo"].notna())].copy()
df_vence["dias"] = df_vence["prazo"].apply(dias_para_prazo)
df_vence_7 = df_vence[(df_vence["dias"] >= 0) & (df_vence["dias"] <= 7)]
if not df_vence_7.empty:
    nomes = ", ".join([f"#{int(r['numero'])} ({r['responsavel']})" for _, r in df_vence_7.iterrows()])
    st.warning(f"🔔 **{len(df_vence_7)} ação(ões) vencem nos próximos 7 dias:** {nomes}")

st.divider()

# TABS
aba_acoes, aba_dashboard, aba_responsavel, aba_ata, aba_nova, aba_export = st.tabs([
    "📋 Ações", "📊 Dashboard", "👤 Por Responsável", "📝 Ata", "➕ Nova Ação", "📥 Exportar"
])

# ============================
# ABA AÇÕES
# ============================
with aba_acoes:
    col_f1, col_f2, col_f3 = st.columns(3)
    with col_f1:
        st.markdown("👤 **Responsável**")
        filtro_resp = st.selectbox("", ["Todos"] + RESPONSAVEIS, label_visibility="collapsed", key="f_resp")
    with col_f2:
        st.markdown("📌 **Status**")
        filtro_status = st.selectbox("", ["Todos"] + STATUS_OPCOES, label_visibility="collapsed", key="f_status")
    with col_f3:
        st.markdown("🏷️ **Tipo**")
        filtro_tipo = st.selectbox("", ["Todos"] + TIPOS, label_visibility="collapsed", key="f_tipo")

    col_f4, col_f5, col_f6 = st.columns(3)
    with col_f4:
        st.markdown("📅 **Data entrada — de**")
        data_de = st.date_input("", value=None, label_visibility="collapsed", key="data_de")
    with col_f5:
        st.markdown("📅 **Data entrada — até**")
        data_ate = st.date_input("", value=None, label_visibility="collapsed", key="data_ate")
    with col_f6:
        st.markdown("🔍 **Buscar**")
        busca = st.text_input("", placeholder="Digite para buscar...", label_visibility="collapsed", key="busca")

    df = df_all.copy()
    if filtro_resp != "Todos": df = df[df["responsavel"].str.contains(filtro_resp, case=False, na=False)]
    if filtro_status != "Todos": df = df[df["status"] == filtro_status]
    if filtro_tipo != "Todos": df = df[df["tipo"] == filtro_tipo]
    if data_de: df = df[df["data_entrada"] >= str(data_de)]
    if data_ate: df = df[df["data_entrada"] <= str(data_ate)]
    if busca:
        df = df[df["problema_identificado"].str.contains(busca, case=False, na=False) |
                df["plano_de_acao"].str.contains(busca, case=False, na=False)]

    st.caption(f"Mostrando **{len(df)}** ações")
    st.divider()

    for _, row in df.iterrows():
        prazo_str = limpar(row.get('prazo')) or '-'
        icon = STATUS_ICONS.get(row['status'], "")
        atualizado_por = limpar(row.get('atualizado_por'))
        atualizado_em = formatar_data(row.get('atualizado_em'))
        comentario_atual = limpar(row.get('comentario'))
        observacao_atual = limpar(row.get('observacao'))
        tipo_atual = limpar(row.get('tipo')) or 'Antigo'

        dias = dias_para_prazo(row.get('prazo'))
        badge = ""
        if row['status'] not in ["Concluído", "Cancelado"] and dias is not None:
            if dias < 0: badge = f" 🔴 Vencido há {abs(dias)}d"
            elif dias == 0: badge = " 🟠 Vence hoje!"
            elif dias <= 7: badge = f" 🟡 Vence em {dias}d"

        with st.expander(f"#{row['numero']} — {row['responsavel']} | {icon} {row['status']} | Prazo: {prazo_str}{badge} | 🏷️ {tipo_atual}"):
            st.markdown(f"**Problema:**\n{row['problema_identificado']}")
            st.markdown(f"**Plano de Ação:**\n{row['plano_de_acao']}")
            if observacao_atual: st.caption(f"📝 {observacao_atual}")
            if comentario_atual:
                info_str = f"💬 **Comentário:** {comentario_atual}"
                if atualizado_por and atualizado_em: info_str += f" · Atualizado por **{atualizado_por}** em {atualizado_em}"
                st.info(info_str)
            elif atualizado_por and atualizado_em:
                st.caption(f"🕐 Atualizado por **{atualizado_por}** em {atualizado_em}")

            hist = carregar_historico(row['id'])
            if not hist.empty:
                with st.expander("📜 Ver histórico de alterações"):
                    for _, h in hist.iterrows():
                        dt_h = formatar_data(h.get('atualizado_em'))
                        prazo_h = f" | Novo prazo: {h['novo_prazo']}" if limpar(h.get('novo_prazo')) else ""
                        coment_h = f" | Comentário: {h['comentario']}" if limpar(h.get('comentario')) else ""
                        st.markdown(f"- **{dt_h}** — {h.get('atualizado_por','')} mudou de *{h.get('status_anterior','')}* para *{h.get('status_novo','')}*{prazo_h}{coment_h}")

            if usuario != "Selecione...":
                col_s, col_t, col_c = st.columns([1, 1, 2])
                with col_s:
                    novo_status = st.selectbox("Novo status", STATUS_OPCOES,
                        index=STATUS_OPCOES.index(row["status"]) if row["status"] in STATUS_OPCOES else 0,
                        key=f"status_{row['id']}")
                with col_t:
                    novo_tipo = st.selectbox("Tipo", TIPOS,
                        index=TIPOS.index(tipo_atual) if tipo_atual in TIPOS else 0,
                        key=f"tipo_{row['id']}")
                with col_c:
                    label_c = "Comentário (obrigatório) *" if novo_status == "Atrasado" else "Comentário (opcional)"
                    comentario = st.text_input(label_c, value="", key=f"comentario_{row['id']}")

                novo_prazo = None
                mudando_de_atrasado = (row['status'] == "Atrasado" and novo_status != "Atrasado")
                if mudando_de_atrasado:
                    st.warning("⚠️ Ação estava atrasada — informe o novo prazo.")
                    novo_prazo = st.date_input("Novo prazo *", value=None, key=f"prazo_{row['id']}")

                if st.button("💾 Salvar", key=f"salvar_{row['id']}"):
                    status_mudou = novo_status != row['status']
                    tipo_mudou = novo_tipo != tipo_atual
                    comentario_mudou = comentario.strip() != ""
                    if novo_status == "Atrasado" and not comentario.strip():
                        st.error("Comentário obrigatório ao marcar como Atrasado.")
                    elif mudando_de_atrasado and not novo_prazo:
                        st.error("Informe o novo prazo.")
                    elif mudando_de_atrasado and not comentario.strip():
                        st.error("Comentário obrigatório ao alterar o prazo.")
                    elif not status_mudou and not comentario_mudou and not tipo_mudou:
                        st.warning("Nenhuma alteração detectada.")
                    else:
                        prazo_salvar = novo_prazo.isoformat() if novo_prazo else None
                        if status_mudou:
                            salvar_historico(row['id'], row['numero'], row['status'], novo_status, comentario, prazo_salvar, usuario)
                        atualizar_registro(row['id'], novo_status, comentario or None, usuario, prazo_salvar, novo_tipo)
                        st.success(f"✅ Salvo em {agora_brasilia().strftime('%d/%m/%Y %H:%M')}")
                        st.rerun()
            else:
                st.warning("Selecione seu nome na barra lateral para editar.")

# ============================
# ABA DASHBOARD
# ============================
with aba_dashboard:
    st.subheader("📊 Dashboard — Análise por Tipo")

    col_d1, col_d2 = st.columns(2)
    with col_d1:
        d_de = st.date_input("De", value=None, key="d_de")
    with col_d2:
        d_ate = st.date_input("Até", value=None, key="d_ate")

    df_dash = df_all.copy()
    df_dash["data_entrada"] = pd.to_datetime(df_dash["data_entrada"], errors="coerce")
    if d_de: df_dash = df_dash[df_dash["data_entrada"] >= pd.Timestamp(d_de)]
    if d_ate: df_dash = df_dash[df_dash["data_entrada"] <= pd.Timestamp(d_ate)]
    df_dash["tipo"] = df_dash["tipo"].fillna("Antigo")

    st.divider()

    # Gráfico 1 — Ações por tipo ao longo do tempo
    st.markdown("#### 📅 Ações por tipo ao longo do tempo")
    df_dash["mes"] = df_dash["data_entrada"].dt.to_period("M").astype(str)
    evo_tipo = df_dash.groupby(["mes", "tipo"]).size().unstack(fill_value=0).reset_index()
    if not evo_tipo.empty:
        st.bar_chart(evo_tipo.set_index("mes"))
    else:
        st.info("Sem dados para o período selecionado.")

    st.divider()

    # Gráfico 2 — Taxa de conclusão por tipo
    st.markdown("#### ✅ Taxa de conclusão por tipo")
    taxa_tipo = df_dash.groupby("tipo").apply(
        lambda x: round(len(x[x["status"] == "Concluído"]) / len(x) * 100, 1) if len(x) > 0 else 0
    ).reset_index(name="Taxa (%)")
    taxa_tipo = taxa_tipo.sort_values("Taxa (%)", ascending=False)
    st.bar_chart(taxa_tipo.set_index("tipo"))

    st.divider()

    # Gráfico 3 — Comparativo de atrasos por tipo
    st.markdown("#### ⚠️ Ações atrasadas por tipo")
    atraso_tipo = df_dash[df_dash["status"] == "Atrasado"].groupby("tipo").size().reset_index(name="Atrasadas")
    if not atraso_tipo.empty:
        st.bar_chart(atraso_tipo.set_index("tipo"))
    else:
        st.info("Nenhuma ação atrasada no período.")

    st.divider()

    # Tabela resumo
    st.markdown("#### 📋 Resumo por tipo")
    resumo_tipo = df_dash.groupby("tipo").agg(
        Total=("status", "count"),
        Concluídas=("status", lambda x: (x == "Concluído").sum()),
        Atrasadas=("status", lambda x: (x == "Atrasado").sum()),
        Em_andamento=("status", lambda x: (x == "Em andamento").sum()),
    ).reset_index()
    resumo_tipo["Taxa (%)"] = (resumo_tipo["Concluídas"] / resumo_tipo["Total"] * 100).round(1)
    resumo_tipo = resumo_tipo.sort_values("Total", ascending=False)
    st.dataframe(resumo_tipo, use_container_width=True, hide_index=True)

# ============================
# ABA POR RESPONSÁVEL
# ============================
with aba_responsavel:
    st.subheader("👤 Visão por Responsável")
    resumo = df_all.groupby(["responsavel", "status"]).size().unstack(fill_value=0).reset_index()
    for col in STATUS_OPCOES:
        if col not in resumo.columns: resumo[col] = 0
    resumo["Total"] = resumo[STATUS_OPCOES].sum(axis=1)
    resumo = resumo.sort_values("Total", ascending=False).rename(columns={"responsavel": "Responsável"})
    st.dataframe(resumo[["Responsável", "Em andamento", "Atrasado", "Concluído", "Cancelado", "Total"]],
                 use_container_width=True, hide_index=True)

    st.divider()
    resp_sel = st.selectbox("Ver ações de:", RESPONSAVEIS, key="resp_vis")
    df_resp = df_all[df_all["responsavel"].str.contains(resp_sel, case=False, na=False)]
    st.caption(f"{len(df_resp)} ações de {resp_sel}")
    for _, row in df_resp.iterrows():
        icon = STATUS_ICONS.get(row['status'], "")
        st.markdown(f"- **#{row['numero']}** {icon} {row['status']} | Prazo: {limpar(row.get('prazo')) or '-'} | {str(row['problema_identificado'])[:80]}...")

# ============================
# ABA ATA
# ============================
with aba_ata:
    st.subheader("📝 Gerador de Ata")

    col_a1, col_a2 = st.columns(2)
    with col_a1:
        ata_status = st.multiselect("Incluir status:", STATUS_OPCOES, default=["Em andamento", "Atrasado"])
    with col_a2:
        ata_resp = st.multiselect("Filtrar por responsável (opcional):", RESPONSAVEIS)

    titulo_ata = st.text_input("Título da reunião", value=f"Fórum de Portabilidade — {agora_brasilia().strftime('%d/%m/%Y')}")

    df_ata = df_all[df_all["status"].isin(ata_status)]
    if ata_resp:
        df_ata = df_ata[df_ata["responsavel"].isin(ata_resp)]

    col_b1, col_b2 = st.columns(2)

    with col_b1:
        if st.button("📝 Gerar texto da ata", type="primary"):
            if df_ata.empty:
                st.info("Nenhuma ação encontrada.")
            else:
                ata = f"**ATA — {titulo_ata}**\n\nTotal de ações: {len(df_ata)}\n\n---\n\n"
                for _, row in df_ata.iterrows():
                    icon = STATUS_ICONS.get(row['status'], "")
                    prazo = limpar(row.get('prazo')) or 'Sem prazo'
                    tipo = limpar(row.get('tipo')) or '-'
                    comentario = limpar(row.get('comentario'))
                    atualizado_por = limpar(row.get('atualizado_por'))
                    ata += f"**#{row['numero']} — {row['responsavel']}** | {icon} {row['status']} | Prazo: {prazo} | Tipo: {tipo}\n"
                    ata += f"📌 {row['problema_identificado']}\n🎯 {row['plano_de_acao']}\n"
                    if comentario:
                        ata += f"💬 {comentario}"
                        if atualizado_por: ata += f" ({atualizado_por})"
                        ata += "\n"
                    ata += "\n"
                st.text_area("Copie o texto:", value=ata, height=400)

    with col_b2:
        if st.button("📊 Exportar ata em PPT"):
            if df_ata.empty:
                st.info("Nenhuma ação encontrada.")
            else:
                with st.spinner("Gerando PPT..."):
                    buf = gerar_ppt_ata(df_ata, titulo_ata)
                nome_arquivo = f"ata_{agora_brasilia().strftime('%Y%m%d_%H%M')}.pptx"
                st.download_button(
                    label="⬇️ Baixar PPT da Ata",
                    data=buf,
                    file_name=nome_arquivo,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

# ============================
# ABA NOVA AÇÃO
# ============================
with aba_nova:
    if usuario == "Selecione...":
        st.warning("Selecione seu nome na barra lateral para cadastrar uma nova ação.")
    else:
        st.subheader("➕ Cadastrar Nova Ação")
        proximo_numero = int(df_all["numero"].max()) + 1 if len(df_all) > 0 else 1
        st.caption(f"Número da ação: **#{proximo_numero}**")
        problema = st.text_area("Problema Identificado *", height=100)
        plano = st.text_area("Plano de Ação *", height=100)
        col_n1, col_n2 = st.columns(2)
        with col_n1:
            responsaveis_acao = st.multiselect("Responsável(is) *", RESPONSAVEIS, key="resp_nova")
        with col_n2:
            tipo_novo = st.selectbox("Tipo *", TIPOS, key="tipo_nova")
        prazo = st.date_input("Prazo", value=None, key="prazo_nova")
        observacao = st.text_input("Observação (opcional)")

        if st.button("✅ Cadastrar Ação", type="primary"):
            if not problema or not plano or not responsaveis_acao:
                st.error("Preencha os campos obrigatórios.")
            else:
                responsavel_texto = ", ".join(responsaveis_acao)
                dados = {
                    "numero": proximo_numero, "origem": "Fórum de Portabilidade",
                    "data_entrada": agora_brasilia().strftime("%Y-%m-%d"),
                    "problema_identificado": problema, "plano_de_acao": plano,
                    "responsavel": responsavel_texto, "tipo": tipo_novo,
                    "prazo": prazo.isoformat() if prazo else None,
                    "data_finalizacao": None, "status": "Em andamento",
                    "dias_atraso": None, "observacao": observacao or None,
                    "comentario": None, "atualizado_por": usuario,
                    "atualizado_em": agora_brasilia().isoformat(),
                }
                cadastrar_acao(dados)
                st.success(f"✅ Ação #{proximo_numero} cadastrada!")
                st.rerun()

# ============================
# ABA EXPORTAR
# ============================
with aba_export:
    st.subheader("📥 Exportar dados para Excel")
    df_export = df_all.copy().rename(columns={
        "numero": "Número", "origem": "Origem", "data_entrada": "Data Entrada",
        "problema_identificado": "Problema Identificado", "plano_de_acao": "Plano de Ação",
        "responsavel": "Responsável", "prazo": "Prazo", "data_finalizacao": "Data Finalização",
        "status": "Status", "dias_atraso": "Dias Atraso", "observacao": "Observação",
        "tipo": "Tipo", "comentario": "Último Comentário",
        "atualizado_por": "Atualizado Por", "atualizado_em": "Atualizado Em",
    }).drop(columns=["id"], errors="ignore")

    buffer = BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        df_export.to_excel(writer, index=False, sheet_name="Plano de Ação")
    buffer.seek(0)
    st.download_button(
        label="⬇️ Baixar Excel",
        data=buffer,
        file_name=f"plano_acao_{agora_brasilia().strftime('%Y%m%d_%H%M')}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
