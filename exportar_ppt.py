from datetime import datetime
2
import math
3
 
4
from pptx import Presentation
5
from pptx.util import Inches, Pt
6
from pptx.dml.color import RGBColor
7
from pptx.enum.text import PP_ALIGN
8
 
9
COR_VERMELHO = RGBColor(192, 0, 0)
10
COR_VERDE = RGBColor(112, 173, 71)
11
COR_LARANJA = RGBColor(244, 177, 131)
12
COR_CINZA = RGBColor(127, 127, 127)
13
COR_CINZA_CLARO = RGBColor(230, 230, 230)
14
 
15
 
16
def gerar_ppt(template_path, df_filtrado, filtros_texto, saida_path):
17
 
18
prs = Presentation()
19
 
20
# ==========================
21
# CAPA
22
# ==========================
23
 
24
slide = prs.slides.add_slide(prs.slide_layouts[6])
25
 
26
titulo = slide.shapes.add_textbox(
27
Inches(1),
28
Inches(1.5),
29
Inches(7),
30
Inches(1)
31
)
32
 
33
titulo.text_frame.text = "Plano de Ação\nPortabilidade"
34
 
35
for p in titulo.text_frame.paragraphs:
36
p.alignment = PP_ALIGN.LEFT
37
for run in p.runs:
38
run.font.size = Pt(28)
39
run.font.bold = True
40
run.font.name = "AMX"
41
run.font.color.rgb = COR_VERMELHO
42
 
43
sub = slide.shapes.add_textbox(
44
Inches(1),
45
Inches(3),
46
Inches(5),
47
Inches(0.5)
48
)
49
 
50
sub.text_frame.text = (
51
f"Gerado em {datetime.now().strftime('%d/%m/%Y %H:%M')}"
52
)
53
 
54
# ==========================
55
# RESUMO
56
# ==========================
57
 
58
slide = prs.slides.add_slide(prs.slide_layouts[6])
59
 
60
titulo = slide.shapes.add_textbox(
61
Inches(0.5),
62
Inches(0.3),
63
Inches(4),
64
Inches(0.5)
65
)
66
 
67
titulo.text_frame.text = "Resumo Executivo"
68
 
69
total = len(df_filtrado)
70
 
71
concluidas = len(
72
df_filtrado[
73
df_filtrado["status_exibicao"] == "Concluído"
74
]
75
)
76
 
77
atrasadas = len(
78
df_filtrado[
79
df_filtrado["status_exibicao"] == "Atrasado"
80
]
81
)
82
 
83
andamento = len(
84
df_filtrado[
85
df_filtrado["status_exibicao"] == "Em andamento"
86
]
87
)
88
 
89
estagnadas = int(
90
df_filtrado["estagnada"].sum()
91
) if len(df_filtrado) > 0 else 0
92
 
93
taxa = (
94
round(concluidas / total * 100, 1)
95
if total > 0 else 0
96
)
97
 
98
cards = [
99
("Total", total, COR_CINZA),
100
("Concluídas", concluidas, COR_VERDE),
101
("Atrasadas", atrasadas, COR_VERMELHO),
102
("Em andamento", andamento, COR_LARANJA),
103
("Estagnadas", estagnadas, COR_VERMELHO),
104
("Taxa", f"{taxa}%", COR_VERMELHO),
105
]
106
 
107
for i, (titulo_card, valor, cor) in enumerate(cards):
108
 
109
x = 0.4 + (i * 1.6)
110
 
111
card = slide.shapes.add_textbox(
112
Inches(x),
113
Inches(1),
114
Inches(1.4),
115
Inches(0.7)
116
)
117
 
118
card.fill.solid()
119
card.fill.fore_color.rgb = cor
120
 
121
tf = card.text_frame
122
tf.text = f"{valor}\n{titulo_card}"
123
 
124
# ==========================
125
# TABELAS POR STATUS
126
# ==========================
127
 
128
status_ordem = [
129
"Atrasado",
130
"Em andamento",
131
"Concluído"
132
]
133
 
134
linhas_por_slide = 8
135
 
136
for status in status_ordem:
137
 
138
df_status = df_filtrado[
139
df_filtrado["status_exibicao"] == status
140
]
141
 
142
if len(df_status) == 0:
143
continue
144
 
145
paginas = math.ceil(
146
len(df_status) / linhas_por_slide
147
)
148
 
149
for pagina in range(paginas):
150
 
151
slide = prs.slides.add_slide(
152
prs.slide_layouts[6]
153
)
154
 
155
titulo = slide.shapes.add_textbox(
156
Inches(0.3),
157
Inches(0.2),
158
Inches(6),
159
Inches(0.4)
160
)
161
 
162
titulo.text_frame.text = (
163
f"Ações {status} ({pagina+1}/{paginas})"
164
)
165
 
166
inicio = pagina * linhas_por_slide
167
fim = inicio + linhas_por_slide
168
 
169
df_pagina = df_status.iloc[inicio:fim]
170
 
171
tabela = slide.shapes.add_table(
172
len(df_pagina) + 1,
173
6,
174
Inches(0.15),
175
Inches(0.8),
176
Inches(9.5),
177
Inches(4.5)
178
).table
179
 
180
tabela.columns[0].width = Inches(0.7)
181
tabela.columns[1].width = Inches(1.3)
182
tabela.columns[2].width = Inches(2.2)
183
tabela.columns[3].width = Inches(2.2)
184
tabela.columns[4].width = Inches(1.1)
185
tabela.columns[5].width = Inches(1.5)
186
 
187
headers = [
188
"Número",
189
"Tipo",
190
"Problema",
191
"Plano de Ação",
192
"Status",
193
"Última Atualização"
194
]
195
 
196
for c, h in enumerate(headers):
197
 
198
cell = tabela.cell(0, c)
199
cell.text = h
200
 
201
cell.fill.solid()
202
cell.fill.fore_color.rgb = COR_VERMELHO
203
 
204
for p in cell.text_frame.paragraphs:
205
 
206
p.alignment = PP_ALIGN.CENTER
207
 
208
for run in p.runs:
209
run.font.bold = True
210
run.font.size = Pt(10)
211
run.font.name = "AMX"
212
run.font.color.rgb = RGBColor(
213
255,
214
255,
215
255
216
)
217
 
218
for r, (_, row) in enumerate(
219
df_pagina.iterrows(),
220
start=1
221
):
222
 
223
valores = [
224
str(row["numero"]),
225
str(row["tipo"]),
226
str(row["problema_identificado"])[:120],
227
str(row["plano_de_acao"])[:120],
228
str(row["status_exibicao"]),
229
str(row["atualizado_em_fmt"])
230
]
231
 
232
for c, valor in enumerate(valores):
233
 
234
cell = tabela.cell(r, c)
235
 
236
cell.text = valor
237
 
238
cell.fill.solid()
239
cell.fill.fore_color.rgb = COR_CINZA_CLARO
240
 
241
prs.save(saida_path)
242
 
243
return saida_path
